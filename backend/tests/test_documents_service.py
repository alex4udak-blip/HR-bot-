"""
Comprehensive unit tests for document parsing service.
Tests all document parsing functions, file upload handling, and error scenarios.
"""
import pytest
import asyncio
import io
import json
import tempfile
import zipfile
from pathlib import Path
from datetime import datetime
from unittest.mock import Mock, MagicMock, AsyncMock, patch, mock_open
from typing import List

import pdfplumber
from docx import Document as DocxDocument
from pptx import Presentation
import openpyxl
import pandas as pd
from bs4 import BeautifulSoup

from api.services.documents import (
    DocumentParser,
    DocumentParseResult,
    MAX_FILE_SIZE,
    MAX_PDF_PAGES,
    MAX_ARCHIVE_FILES,
    PARSE_TIMEOUT
)


# ============================================================================
# FIXTURES
# ============================================================================

@pytest.fixture
def parser():
    """Create a DocumentParser instance with mocked Anthropic client."""
    with patch('api.services.documents.anthropic.Anthropic') as mock_anthropic:
        mock_client = Mock()
        mock_anthropic.return_value = mock_client
        parser = DocumentParser()
        parser.client = mock_client
        yield parser


@pytest.fixture
def sample_pdf_bytes():
    """Create a minimal valid PDF for testing."""
    # Minimal PDF structure
    pdf_content = b"""%PDF-1.4
1 0 obj
<< /Type /Catalog /Pages 2 0 R >>
endobj
2 0 obj
<< /Type /Pages /Kids [3 0 R] /Count 1 >>
endobj
3 0 obj
<< /Type /Page /Parent 2 0 R /Resources 4 0 R /MediaBox [0 0 612 792] /Contents 5 0 R >>
endobj
4 0 obj
<< /Font << /F1 << /Type /Font /Subtype /Type1 /BaseFont /Helvetica >> >> >>
endobj
5 0 obj
<< /Length 44 >>
stream
BT
/F1 12 Tf
100 700 Td
(Test PDF) Tj
ET
endstream
endobj
xref
0 6
0000000000 65535 f
0000000009 00000 n
0000000058 00000 n
0000000115 00000 n
0000000214 00000 n
0000000301 00000 n
trailer
<< /Size 6 /Root 1 0 R >>
startxref
395
%%EOF"""
    return pdf_content


@pytest.fixture
def sample_text_bytes():
    """Create sample text content."""
    return b"This is a test document.\nWith multiple lines.\nAnd some content."


@pytest.fixture
def sample_json_bytes():
    """Create sample JSON content."""
    data = {
        "name": "John Doe",
        "position": "Software Engineer",
        "skills": ["Python", "JavaScript", "SQL"]
    }
    return json.dumps(data, indent=2).encode('utf-8')


@pytest.fixture
def sample_html_bytes():
    """Create sample HTML content."""
    html = """
    <!DOCTYPE html>
    <html>
    <head><title>Test Page</title></head>
    <body>
        <h1>Resume</h1>
        <p>Name: John Doe</p>
        <script>console.log('test');</script>
        <style>body { color: black; }</style>
    </body>
    </html>
    """
    return html.encode('utf-8')


@pytest.fixture
def sample_csv_bytes():
    """Create sample CSV content."""
    csv = "Name,Position,Department\nJohn Doe,Engineer,IT\nJane Smith,Manager,HR"
    return csv.encode('utf-8')


@pytest.fixture
def sample_xml_bytes():
    """Create sample XML content."""
    xml = """<?xml version="1.0"?>
    <employee>
        <name>John Doe</name>
        <position>Engineer</position>
        <department>IT</department>
    </employee>
    """
    return xml.encode('utf-8')


@pytest.fixture
def sample_email_bytes():
    """Create sample email content."""
    email = """From: sender@example.com
To: recipient@example.com
Subject: Test Email
Date: Mon, 1 Jan 2024 12:00:00 +0000

This is the email body.
It contains important information.
"""
    return email.encode('utf-8')


# ============================================================================
# TEST DOCUMENTPARSERESULT
# ============================================================================

class TestDocumentParseResult:
    """Tests for DocumentParseResult class."""

    def test_create_result_with_defaults(self):
        """Test creating result with default values."""
        result = DocumentParseResult()

        assert result.content == ""
        assert result.status == "parsed"
        assert result.metadata == {}
        assert result.error is None

    def test_create_result_with_content(self):
        """Test creating result with content."""
        result = DocumentParseResult(
            content="Test content",
            status="parsed",
            metadata={"filename": "test.txt"}
        )

        assert result.content == "Test content"
        assert result.status == "parsed"
        assert result.metadata["filename"] == "test.txt"
        assert result.error is None

    def test_create_result_with_error(self):
        """Test creating result with error."""
        result = DocumentParseResult(
            status="failed",
            error="Parse error",
            metadata={"filename": "test.pdf"}
        )

        assert result.content == ""
        assert result.status == "failed"
        assert result.error == "Parse error"
        assert result.metadata["filename"] == "test.pdf"

    def test_create_result_partial_status(self):
        """Test creating result with partial status."""
        result = DocumentParseResult(
            content="Some content",
            status="partial",
            error="Could not parse all pages"
        )

        assert result.status == "partial"
        assert result.content == "Some content"
        assert result.error is not None


# ============================================================================
# TEST MAIN PARSE ENTRY POINT
# ============================================================================

class TestDocumentParserMain:
    """Tests for main parse() entry point."""

    @pytest.mark.asyncio
    async def test_parse_file_too_large(self, parser):
        """Test that files exceeding MAX_FILE_SIZE are rejected."""
        large_file = b"x" * (MAX_FILE_SIZE + 1)

        result = await parser.parse(large_file, "large.pdf")

        assert result.status == "failed"
        assert "File too large" in result.error
        assert result.metadata["filename"] == "large.pdf"
        assert result.metadata["file_size"] == len(large_file)

    @pytest.mark.asyncio
    async def test_parse_unsupported_file_type(self, parser):
        """Test that unsupported file types return error."""
        result = await parser.parse(b"test", "document.unsupported")

        assert result.status == "failed"
        assert "Unsupported file type" in result.error
        assert result.metadata["file_type"] == "unsupported"

    @pytest.mark.asyncio
    async def test_parse_adds_metadata(self, parser, sample_text_bytes):
        """Test that parse adds standard metadata."""
        result = await parser.parse(sample_text_bytes, "test.txt")

        assert result.metadata["filename"] == "test.txt"
        assert result.metadata["file_type"] == "txt"
        assert result.metadata["file_size"] == len(sample_text_bytes)
        assert "parsed_at" in result.metadata

    @pytest.mark.asyncio
    async def test_parse_timeout_handling(self, parser):
        """Test that parsing timeout is handled properly."""
        async def slow_parser(*args, **kwargs):
            await asyncio.sleep(PARSE_TIMEOUT + 1)
            return DocumentParseResult()

        with patch.object(parser, '_parse_text', side_effect=slow_parser):
            result = await parser.parse(b"test", "test.txt")

        assert result.status == "failed"
        assert "timeout" in result.error.lower()

    @pytest.mark.asyncio
    async def test_parse_exception_handling(self, parser):
        """Test that parsing exceptions are caught and reported."""
        with patch.object(parser, '_parse_text', side_effect=Exception("Test error")):
            result = await parser.parse(b"test", "test.txt")

        assert result.status == "failed"
        assert "Test error" in result.error


# ============================================================================
# TEST PDF PARSING
# ============================================================================

class TestPDFParsing:
    """Tests for PDF parsing functionality."""

    @pytest.mark.asyncio
    async def test_parse_pdf_success(self, parser):
        """Test successful PDF parsing."""
        mock_page = Mock()
        mock_page.extract_text.return_value = "Page content"
        mock_page.extract_tables.return_value = []

        mock_pdf = Mock()
        mock_pdf.pages = [mock_page]
        mock_pdf.__enter__ = Mock(return_value=mock_pdf)
        mock_pdf.__exit__ = Mock(return_value=False)

        with patch('pdfplumber.open', return_value=mock_pdf):
            result = await parser.parse(b"pdf content", "test.pdf")

        assert result.status == "parsed"
        assert "Page content" in result.content
        assert result.metadata["pages_count"] == 1
        assert result.metadata["pages_parsed"] == 1

    @pytest.mark.asyncio
    async def test_parse_pdf_with_tables(self, parser):
        """Test PDF parsing with tables."""
        mock_page = Mock()
        mock_page.extract_text.return_value = "Text content"
        mock_page.extract_tables.return_value = [
            [["Header1", "Header2"], ["Cell1", "Cell2"]]
        ]

        mock_pdf = Mock()
        mock_pdf.pages = [mock_page]
        mock_pdf.__enter__ = Mock(return_value=mock_pdf)
        mock_pdf.__exit__ = Mock(return_value=False)

        with patch('pdfplumber.open', return_value=mock_pdf):
            result = await parser.parse(b"pdf content", "test.pdf")

        assert result.status == "parsed"
        assert "[Table 1]" in result.content
        assert "Header1 | Header2" in result.content
        assert result.metadata["tables_count"] == 1

    @pytest.mark.asyncio
    async def test_parse_pdf_respects_max_pages(self, parser):
        """Test that PDF parsing respects MAX_PDF_PAGES limit."""
        mock_pages = []
        for i in range(MAX_PDF_PAGES + 10):
            mock_page = Mock()
            mock_page.extract_text.return_value = f"Page {i + 1}"
            mock_page.extract_tables.return_value = []
            mock_pages.append(mock_page)

        mock_pdf = Mock()
        mock_pdf.pages = mock_pages
        mock_pdf.__enter__ = Mock(return_value=mock_pdf)
        mock_pdf.__exit__ = Mock(return_value=False)

        with patch('pdfplumber.open', return_value=mock_pdf):
            result = await parser.parse(b"pdf content", "large.pdf")

        assert result.status == "parsed"
        assert result.metadata["pages_count"] == MAX_PDF_PAGES + 10
        assert result.metadata["pages_parsed"] == MAX_PDF_PAGES
        assert result.metadata["truncated"] is True

    @pytest.mark.asyncio
    async def test_parse_pdf_empty_content(self, parser):
        """Test PDF with no extractable text."""
        mock_page = Mock()
        mock_page.extract_text.return_value = ""
        mock_page.extract_tables.return_value = []

        mock_pdf = Mock()
        mock_pdf.pages = [mock_page]
        mock_pdf.__enter__ = Mock(return_value=mock_pdf)
        mock_pdf.__exit__ = Mock(return_value=False)

        with patch('pdfplumber.open', return_value=mock_pdf):
            result = await parser.parse(b"pdf content", "empty.pdf")

        assert result.status == "partial"
        assert result.content == ""

    def test_table_to_text(self, parser):
        """Test table to text conversion."""
        table = [
            ["Name", "Age", "City"],
            ["John", "30", "NYC"],
            ["Jane", "25", "LA"]
        ]

        text = parser._table_to_text(table)

        assert "Name | Age | City" in text
        assert "John | 30 | NYC" in text
        assert "Jane | 25 | LA" in text

    def test_table_to_text_with_none_values(self, parser):
        """Test table conversion with None values."""
        table = [
            ["Name", None, "City"],
            [None, "30", None]
        ]

        text = parser._table_to_text(table)

        assert "Name |  | City" in text
        assert " | 30 | " in text


# ============================================================================
# TEST DOCX PARSING
# ============================================================================

class TestDOCXParsing:
    """Tests for DOCX parsing functionality."""

    @pytest.mark.asyncio
    async def test_parse_docx_success(self, parser):
        """Test successful DOCX parsing."""
        mock_para1 = Mock()
        mock_para1.text = "First paragraph"
        mock_para2 = Mock()
        mock_para2.text = "Second paragraph"

        mock_doc = Mock()
        mock_doc.paragraphs = [mock_para1, mock_para2]
        mock_doc.tables = []

        with patch('api.services.documents.DocxDocument', return_value=mock_doc):
            result = await parser.parse(b"docx content", "test.docx")

        assert result.status == "parsed"
        assert "First paragraph" in result.content
        assert "Second paragraph" in result.content
        assert result.metadata["paragraphs_count"] == 2

    @pytest.mark.asyncio
    async def test_parse_docx_with_tables(self, parser):
        """Test DOCX parsing with tables."""
        mock_para = Mock()
        mock_para.text = "Text content"

        # Create mock table
        mock_cell1 = Mock()
        mock_cell1.text = "Header1"
        mock_cell2 = Mock()
        mock_cell2.text = "Header2"
        mock_row = Mock()
        mock_row.cells = [mock_cell1, mock_cell2]

        mock_table = Mock()
        mock_table.rows = [mock_row]

        mock_doc = Mock()
        mock_doc.paragraphs = [mock_para]
        mock_doc.tables = [mock_table]

        with patch('api.services.documents.DocxDocument', return_value=mock_doc):
            result = await parser.parse(b"docx content", "test.docx")

        assert result.status == "parsed"
        assert "[Table 1]" in result.content
        assert "Header1 | Header2" in result.content
        assert result.metadata["tables_count"] == 1

    @pytest.mark.asyncio
    async def test_parse_docx_empty(self, parser):
        """Test DOCX with no content."""
        mock_doc = Mock()
        mock_doc.paragraphs = []
        mock_doc.tables = []

        with patch('api.services.documents.DocxDocument', return_value=mock_doc):
            result = await parser.parse(b"docx content", "empty.docx")

        assert result.status == "partial"
        assert result.content == ""

    @pytest.mark.asyncio
    async def test_parse_docx_filters_empty_paragraphs(self, parser):
        """Test that empty paragraphs are filtered out."""
        mock_para1 = Mock()
        mock_para1.text = "Content"
        mock_para2 = Mock()
        mock_para2.text = "   "  # Whitespace only
        mock_para3 = Mock()
        mock_para3.text = ""  # Empty

        mock_doc = Mock()
        mock_doc.paragraphs = [mock_para1, mock_para2, mock_para3]
        mock_doc.tables = []

        with patch('api.services.documents.DocxDocument', return_value=mock_doc):
            result = await parser.parse(b"docx content", "test.docx")

        assert result.status == "parsed"
        assert result.content == "Content"


# ============================================================================
# TEST CONVERSION-BASED PARSING (DOC, ODT, XLS, etc.)
# ============================================================================

class TestConversionParsing:
    """Tests for parsing that requires LibreOffice conversion."""

    @pytest.mark.asyncio
    async def test_convert_doc_to_docx_success(self, parser):
        """Test successful DOC to DOCX conversion."""
        mock_process = AsyncMock()
        mock_process.communicate = AsyncMock(return_value=(b"", b""))

        mock_doc = Mock()
        mock_doc.paragraphs = [Mock(text="Converted content")]
        mock_doc.tables = []

        with patch('asyncio.create_subprocess_exec', return_value=mock_process), \
             patch('builtins.open', mock_open(read_data=b"docx bytes")), \
             patch('os.path.exists', return_value=True), \
             patch('os.remove'), \
             patch('api.services.documents.DocxDocument', return_value=mock_doc):

            result = await parser.parse(b"doc content", "test.doc")

        assert result.status == "parsed"
        assert "Converted content" in result.content
        assert result.metadata.get("converted_from") == ".doc"

    @pytest.mark.asyncio
    async def test_convert_doc_conversion_failed(self, parser):
        """Test DOC conversion failure."""
        mock_process = AsyncMock()
        mock_process.communicate = AsyncMock(return_value=(b"", b""))

        with patch('asyncio.create_subprocess_exec', return_value=mock_process), \
             patch('builtins.open', mock_open()), \
             patch('os.path.exists', return_value=False), \
             patch('os.remove'):

            result = await parser.parse(b"doc content", "test.doc")

        assert result.status == "failed"
        assert "conversion failed" in result.error.lower()


# ============================================================================
# TEST TEXT PARSING
# ============================================================================

class TestTextParsing:
    """Tests for plain text parsing."""

    @pytest.mark.asyncio
    async def test_parse_text_utf8(self, parser, sample_text_bytes):
        """Test parsing UTF-8 text."""
        with patch('chardet.detect', return_value={'encoding': 'utf-8'}):
            result = await parser.parse(sample_text_bytes, "test.txt")

        assert result.status == "parsed"
        assert "test document" in result.content
        assert result.metadata["encoding"] == "utf-8"

    @pytest.mark.asyncio
    async def test_parse_text_different_encoding(self, parser):
        """Test parsing text with different encoding."""
        # Windows-1251 encoded text
        text = "Тестовый документ"
        encoded = text.encode('cp1251')

        with patch('chardet.detect', return_value={'encoding': 'cp1251'}):
            result = await parser.parse(encoded, "test.txt")

        assert result.status == "parsed"
        assert "Тестовый" in result.content

    @pytest.mark.asyncio
    async def test_parse_text_encoding_fallback(self, parser):
        """Test text parsing with encoding detection failure."""
        with patch('chardet.detect', return_value={'encoding': None}):
            result = await parser.parse(b"test content", "test.txt")

        assert result.status == "parsed"
        assert "test content" in result.content

    @pytest.mark.asyncio
    async def test_parse_markdown(self, parser):
        """Test parsing markdown files."""
        markdown = b"# Header\n\n**Bold text**\n\n- List item"

        result = await parser.parse(markdown, "readme.md")

        assert result.status == "parsed"
        assert "# Header" in result.content


# ============================================================================
# TEST HTML PARSING
# ============================================================================

class TestHTMLParsing:
    """Tests for HTML parsing."""

    @pytest.mark.asyncio
    async def test_parse_html_removes_scripts_and_styles(self, parser, sample_html_bytes):
        """Test that HTML parsing removes script and style tags."""
        result = await parser.parse(sample_html_bytes, "test.html")

        assert result.status == "parsed"
        assert "console.log" not in result.content
        assert "color: black" not in result.content
        assert "Resume" in result.content
        assert "John Doe" in result.content

    @pytest.mark.asyncio
    async def test_parse_html_extracts_title(self, parser):
        """Test HTML title extraction."""
        html = b"<html><head><title>My Resume</title></head><body>Content</body></html>"

        result = await parser.parse(html, "test.html")

        assert result.metadata["title"] == "My Resume"

    @pytest.mark.asyncio
    async def test_parse_html_no_title(self, parser):
        """Test HTML without title."""
        html = b"<html><body>Content</body></html>"

        result = await parser.parse(html, "test.html")

        assert result.metadata["title"] is None

    @pytest.mark.asyncio
    async def test_parse_html_cleans_whitespace(self, parser):
        """Test that HTML parsing cleans up excessive whitespace."""
        html = b"""<html><body>
            <p>Line 1</p>


            <p>Line 2</p>
        </body></html>"""

        result = await parser.parse(html, "test.html")

        assert result.status == "parsed"
        # Should not have excessive blank lines
        assert "\n\n\n" not in result.content


# ============================================================================
# TEST SPREADSHEET PARSING
# ============================================================================

class TestSpreadsheetParsing:
    """Tests for Excel and CSV parsing."""

    @pytest.mark.asyncio
    async def test_parse_xlsx_success(self, parser):
        """Test successful XLSX parsing."""
        # Create mock workbook
        mock_cell1 = Mock()
        mock_cell1.value = "Header1"
        mock_cell2 = Mock()
        mock_cell2.value = "Header2"

        mock_sheet = Mock()
        mock_sheet.iter_rows.return_value = [[mock_cell1, mock_cell2]]

        mock_wb = Mock()
        mock_wb.sheetnames = ["Sheet1"]
        mock_wb.__getitem__ = Mock(return_value=mock_sheet)

        with patch('openpyxl.load_workbook', return_value=mock_wb):
            result = await parser.parse(b"xlsx content", "test.xlsx")

        assert result.status == "parsed"
        assert "Sheet1" in result.content
        assert "Header1 | Header2" in result.content
        assert result.metadata["sheets_count"] == 1

    @pytest.mark.asyncio
    async def test_parse_xlsx_multiple_sheets(self, parser):
        """Test XLSX with multiple sheets."""
        mock_cell = Mock()
        mock_cell.value = "Data"

        mock_sheet = Mock()
        mock_sheet.iter_rows.return_value = [[mock_cell]]

        mock_wb = Mock()
        mock_wb.sheetnames = ["Sheet1", "Sheet2"]
        mock_wb.__getitem__ = Mock(return_value=mock_sheet)

        with patch('openpyxl.load_workbook', return_value=mock_wb):
            result = await parser.parse(b"xlsx content", "test.xlsx")

        assert result.status == "parsed"
        assert "Sheet1" in result.content
        assert "Sheet2" in result.content
        assert result.metadata["sheets_count"] == 2

    @pytest.mark.asyncio
    async def test_parse_csv_success(self, parser, sample_csv_bytes):
        """Test successful CSV parsing."""
        result = await parser.parse(sample_csv_bytes, "test.csv")

        assert result.status == "parsed"
        assert "John Doe" in result.content
        assert "Engineer" in result.content
        assert "Name" in result.metadata["columns"]
        assert result.metadata["rows_count"] == 2

    @pytest.mark.asyncio
    async def test_parse_csv_malformed(self, parser):
        """Test CSV parsing with malformed data."""
        malformed_csv = b"Name,Age\nJohn,30,extra\nJane"

        result = await parser.parse(malformed_csv, "test.csv")

        # Should still return some content even if parsing failed
        assert result.status in ["parsed", "partial"]


# ============================================================================
# TEST PRESENTATION PARSING
# ============================================================================

class TestPresentationParsing:
    """Tests for PowerPoint parsing."""

    @pytest.mark.asyncio
    async def test_parse_pptx_success(self, parser):
        """Test successful PPTX parsing."""
        # Create mock shape with text
        mock_shape1 = Mock()
        mock_shape1.text = "Slide title"
        mock_shape2 = Mock()
        mock_shape2.text = "Slide content"

        mock_slide = Mock()
        mock_slide.shapes = [mock_shape1, mock_shape2]

        mock_prs = Mock()
        mock_prs.slides = [mock_slide]

        with patch('api.services.documents.Presentation', return_value=mock_prs):
            result = await parser.parse(b"pptx content", "test.pptx")

        assert result.status == "parsed"
        assert "Slide 1" in result.content
        assert "Slide title" in result.content
        assert "Slide content" in result.content
        assert result.metadata["slides_count"] == 1

    @pytest.mark.asyncio
    async def test_parse_pptx_filters_empty_shapes(self, parser):
        """Test that empty shapes are filtered out."""
        mock_shape1 = Mock()
        mock_shape1.text = "Content"
        mock_shape2 = Mock()
        mock_shape2.text = "   "  # Whitespace only

        mock_slide = Mock()
        mock_slide.shapes = [mock_shape1, mock_shape2]

        mock_prs = Mock()
        mock_prs.slides = [mock_slide]

        with patch('api.services.documents.Presentation', return_value=mock_prs):
            result = await parser.parse(b"pptx content", "test.pptx")

        assert result.status == "parsed"
        assert "Content" in result.content
        # Should not have extra whitespace from empty shape

    @pytest.mark.asyncio
    async def test_parse_pptx_no_text_content(self, parser):
        """Test PPTX with no text content."""
        mock_slide = Mock()
        mock_slide.shapes = []

        mock_prs = Mock()
        mock_prs.slides = [mock_slide]

        with patch('api.services.documents.Presentation', return_value=mock_prs):
            result = await parser.parse(b"pptx content", "empty.pptx")

        assert result.status == "partial"


# ============================================================================
# TEST RTF PARSING
# ============================================================================

class TestRTFParsing:
    """Tests for RTF parsing."""

    @pytest.mark.asyncio
    async def test_parse_rtf_success(self, parser):
        """Test successful RTF parsing."""
        rtf_content = b"{\\rtf1\\ansi\\deff0 {\\fonttbl {\\f0 Times New Roman;}}\\f0\\fs24 Test RTF content.}"

        with patch('api.services.documents.rtf_to_text', return_value="Test RTF content."):
            result = await parser.parse(rtf_content, "test.rtf")

        assert result.status == "parsed"
        assert "Test RTF content" in result.content

    @pytest.mark.asyncio
    async def test_parse_rtf_encoding_fallback(self, parser):
        """Test RTF parsing tries multiple encodings."""
        rtf_content = b"{\\rtf1 content}"

        # Mock rtf_to_text to fail on first encoding, succeed on second
        call_count = 0
        def mock_rtf_to_text(content):
            nonlocal call_count
            call_count += 1
            if call_count == 1:
                raise UnicodeDecodeError('utf-8', b'', 0, 1, 'test')
            return "Converted content"

        with patch('api.services.documents.rtf_to_text', side_effect=mock_rtf_to_text):
            result = await parser.parse(rtf_content, "test.rtf")

        # Should eventually succeed with different encoding
        assert result.status in ["parsed", "failed"]

    @pytest.mark.asyncio
    async def test_parse_rtf_decode_failure(self, parser):
        """Test RTF parsing when all encodings fail."""
        rtf_content = b"\xff\xfe invalid rtf"

        with patch('api.services.documents.rtf_to_text', side_effect=Exception("Decode error")):
            result = await parser.parse(rtf_content, "test.rtf")

        assert result.status == "failed"
        assert "Could not decode RTF" in result.error


# ============================================================================
# TEST IMAGE OCR
# ============================================================================

class TestImageOCR:
    """Tests for image OCR with Claude Vision."""

    @pytest.mark.asyncio
    async def test_ocr_with_vision_success(self, parser):
        """Test successful OCR with Claude Vision."""
        image_bytes = b"fake image data"

        # Mock Anthropic API response
        mock_content = Mock()
        mock_content.text = "Extracted text from image"
        mock_response = Mock()
        mock_response.content = [mock_content]

        parser.client.messages.create = Mock(return_value=mock_response)

        result = await parser.parse(image_bytes, "test.jpg")

        assert result.status == "parsed"
        assert "Extracted text from image" in result.content
        assert result.metadata["ocr_method"] == "claude_vision"

    @pytest.mark.asyncio
    async def test_ocr_determines_correct_media_type(self, parser):
        """Test that OCR uses correct media type for different image formats."""
        image_bytes = b"fake image data"

        mock_content = Mock()
        mock_content.text = "Text"
        mock_response = Mock()
        mock_response.content = [mock_content]

        parser.client.messages.create = Mock(return_value=mock_response)

        # Test different formats
        formats = {
            "test.jpg": "image/jpeg",
            "test.jpeg": "image/jpeg",
            "test.png": "image/png",
            "test.gif": "image/gif",
            "test.webp": "image/webp"
        }

        for filename, expected_type in formats.items():
            await parser.parse(image_bytes, filename)

            # Check the call was made with correct media type
            call_args = parser.client.messages.create.call_args
            content = call_args[1]["messages"][0]["content"]
            image_content = [c for c in content if c.get("type") == "image"][0]
            assert image_content["source"]["media_type"] == expected_type

    @pytest.mark.asyncio
    async def test_ocr_api_error_handling(self, parser):
        """Test OCR error handling when API fails."""
        image_bytes = b"fake image data"

        parser.client.messages.create = Mock(side_effect=Exception("API error"))

        result = await parser.parse(image_bytes, "test.jpg")

        assert result.status == "failed"
        assert "OCR failed" in result.error

    @pytest.mark.asyncio
    async def test_convert_heic_to_jpeg(self, parser):
        """Test HEIC conversion to JPEG before OCR."""
        # Skip test if pillow_heif is not installed (optional dependency)
        pytest.importorskip('pillow_heif')

        heic_bytes = b"fake heic data"

        # Mock PIL Image
        mock_img = Mock()
        mock_img.convert.return_value.save = Mock()

        mock_content = Mock()
        mock_content.text = "Extracted text"
        mock_response = Mock()
        mock_response.content = [mock_content]

        with patch('PIL.Image.open', return_value=mock_img), \
             patch('pillow_heif.register_heif_opener'):

            parser.client.messages.create = Mock(return_value=mock_response)
            result = await parser.parse(heic_bytes, "test.heic")

        assert result.status == "parsed"


# ============================================================================
# TEST ARCHIVE PARSING
# ============================================================================

class TestArchiveParsing:
    """Tests for ZIP, RAR, and 7z archive parsing."""

    @pytest.mark.asyncio
    async def test_parse_zip_success(self, parser):
        """Test successful ZIP extraction and parsing."""
        # Create a mock ZIP file
        zip_buffer = io.BytesIO()
        with zipfile.ZipFile(zip_buffer, 'w') as zf:
            zf.writestr("document.txt", "Text content")

        zip_bytes = zip_buffer.getvalue()

        result = await parser.parse(zip_bytes, "archive.zip")

        assert result.status == "parsed"
        assert "document.txt" in result.content
        assert "Text content" in result.content
        assert result.metadata["archive_type"] == "zip"
        assert len(result.metadata["extracted_files"]) == 1

    @pytest.mark.asyncio
    async def test_parse_zip_respects_max_files(self, parser):
        """Test that ZIP parsing respects MAX_ARCHIVE_FILES limit."""
        zip_buffer = io.BytesIO()
        with zipfile.ZipFile(zip_buffer, 'w') as zf:
            for i in range(MAX_ARCHIVE_FILES + 5):
                zf.writestr(f"file{i}.txt", f"Content {i}")

        zip_bytes = zip_buffer.getvalue()

        result = await parser.parse(zip_bytes, "large.zip")

        # Should only process MAX_ARCHIVE_FILES
        assert len(result.metadata["extracted_files"]) <= MAX_ARCHIVE_FILES

    @pytest.mark.asyncio
    async def test_parse_zip_skips_directories(self, parser):
        """Test that ZIP parsing skips directory entries."""
        zip_buffer = io.BytesIO()
        with zipfile.ZipFile(zip_buffer, 'w') as zf:
            zf.writestr("folder/", "")  # Directory entry
            zf.writestr("folder/file.txt", "Content")

        zip_bytes = zip_buffer.getvalue()

        result = await parser.parse(zip_bytes, "archive.zip")

        # Should only process the file, not the directory
        assert "Content" in result.content

    @pytest.mark.asyncio
    async def test_parse_zip_handles_nested_parsing_errors(self, parser):
        """Test ZIP parsing when nested files fail to parse."""
        zip_buffer = io.BytesIO()
        with zipfile.ZipFile(zip_buffer, 'w') as zf:
            zf.writestr("good.txt", "Good content")
            zf.writestr("bad.unsupported", "Bad content")

        zip_bytes = zip_buffer.getvalue()

        result = await parser.parse(zip_bytes, "archive.zip")

        # Should handle the good file and note the bad one
        assert len(result.metadata["extracted_files"]) == 2
        failed_files = [f for f in result.metadata["extracted_files"] if f["status"] == "failed"]
        assert len(failed_files) > 0

    @pytest.mark.asyncio
    async def test_parse_rar_success(self, parser):
        """Test successful RAR extraction."""
        mock_process = AsyncMock()
        mock_process.communicate = AsyncMock(return_value=(b"", b""))

        with patch('asyncio.create_subprocess_exec', return_value=mock_process), \
             patch('os.makedirs'), \
             patch('os.walk', return_value=[("/extract", [], ["file.txt"])]), \
             patch('builtins.open', mock_open(read_data=b"Test content")), \
             patch('os.path.exists', return_value=True), \
             patch('os.remove'), \
             patch('shutil.rmtree'):

            result = await parser.parse(b"rar content", "archive.rar")

        assert result.metadata["archive_type"] == "rar"

    @pytest.mark.asyncio
    async def test_parse_7z_success(self, parser):
        """Test successful 7z extraction."""
        mock_bio = io.BytesIO(b"Text content")

        mock_archive = Mock()
        mock_archive.readall.return_value = {"file.txt": mock_bio}
        mock_archive.__enter__ = Mock(return_value=mock_archive)
        mock_archive.__exit__ = Mock(return_value=False)

        with patch('py7zr.SevenZipFile', return_value=mock_archive):
            result = await parser.parse(b"7z content", "archive.7z")

        assert result.metadata["archive_type"] == "7z"


# ============================================================================
# TEST JSON AND XML PARSING
# ============================================================================

class TestStructuredDataParsing:
    """Tests for JSON and XML parsing."""

    @pytest.mark.asyncio
    async def test_parse_json_success(self, parser, sample_json_bytes):
        """Test successful JSON parsing."""
        result = await parser.parse(sample_json_bytes, "data.json")

        assert result.status == "parsed"
        assert "John Doe" in result.content
        assert "Software Engineer" in result.content
        assert result.metadata["json_type"] == "dict"

    @pytest.mark.asyncio
    async def test_parse_json_array(self, parser):
        """Test parsing JSON array."""
        json_data = json.dumps([{"name": "John"}, {"name": "Jane"}]).encode('utf-8')

        result = await parser.parse(json_data, "data.json")

        assert result.status == "parsed"
        assert result.metadata["json_type"] == "list"

    @pytest.mark.asyncio
    async def test_parse_json_invalid(self, parser):
        """Test parsing invalid JSON."""
        invalid_json = b"{invalid json content"

        result = await parser.parse(invalid_json, "data.json")

        assert result.status == "partial"
        assert "Invalid JSON" in result.error

    @pytest.mark.asyncio
    async def test_parse_xml_success(self, parser, sample_xml_bytes):
        """Test successful XML parsing."""
        result = await parser.parse(sample_xml_bytes, "data.xml")

        assert result.status == "parsed"
        assert "John Doe" in result.content
        assert "Engineer" in result.content

    @pytest.mark.asyncio
    async def test_parse_yaml(self, parser):
        """Test parsing YAML files (treated as text)."""
        yaml_content = b"name: John Doe\nposition: Engineer\nskills:\n  - Python\n  - SQL"

        result = await parser.parse(yaml_content, "config.yaml")

        assert result.status == "parsed"
        assert "John Doe" in result.content


# ============================================================================
# TEST EMAIL PARSING
# ============================================================================

class TestEmailParsing:
    """Tests for email parsing (.eml and .msg)."""

    @pytest.mark.asyncio
    async def test_parse_eml_success(self, parser, sample_email_bytes):
        """Test successful EML parsing."""
        result = await parser.parse(sample_email_bytes, "message.eml")

        assert result.status == "parsed"
        assert "sender@example.com" in result.content
        assert "recipient@example.com" in result.content
        assert "Test Email" in result.content
        assert result.metadata["subject"] == "Test Email"

    @pytest.mark.asyncio
    async def test_parse_eml_multipart(self, parser):
        """Test parsing multipart email."""
        email_content = b"""From: sender@test.com
To: recipient@test.com
Subject: Multipart Test
MIME-Version: 1.0
Content-Type: multipart/alternative; boundary="boundary"

--boundary
Content-Type: text/plain

Plain text content
--boundary
Content-Type: text/html

<html><body>HTML content</body></html>
--boundary--
"""

        result = await parser.parse(email_content, "message.eml")

        assert result.status == "parsed"
        assert "Plain text content" in result.content

    @pytest.mark.asyncio
    async def test_parse_outlook_msg_success(self, parser):
        """Test successful Outlook MSG parsing."""
        # Skip test if extract_msg is not installed (optional dependency)
        extract_msg = pytest.importorskip('extract_msg')

        mock_msg = Mock()
        mock_msg.sender = "sender@test.com"
        mock_msg.to = "recipient@test.com"
        mock_msg.subject = "Test Subject"
        mock_msg.date = "2024-01-01"
        mock_msg.body = "Email body content"

        with patch('extract_msg.Message', return_value=mock_msg), \
             patch('builtins.open', mock_open()), \
             patch('os.path.exists', return_value=True), \
             patch('os.remove'):

            result = await parser.parse(b"msg content", "message.msg")

        assert result.status == "parsed"
        assert "sender@test.com" in result.content
        assert "Test Subject" in result.content
        assert "Email body content" in result.content

    @pytest.mark.asyncio
    async def test_parse_msg_error_handling(self, parser):
        """Test MSG parsing error handling."""
        # Skip test if extract_msg is not installed (optional dependency)
        extract_msg = pytest.importorskip('extract_msg')

        with patch('extract_msg.Message', side_effect=Exception("Parse error")), \
             patch('builtins.open', mock_open()), \
             patch('os.path.exists', return_value=True), \
             patch('os.remove'):

            result = await parser.parse(b"msg content", "message.msg")

        assert result.status == "failed"
        assert "MSG parsing failed" in result.error


# ============================================================================
# TEST ERROR HANDLING AND EDGE CASES
# ============================================================================

class TestErrorHandling:
    """Tests for error handling and edge cases."""

    @pytest.mark.asyncio
    async def test_empty_file(self, parser):
        """Test parsing empty file."""
        result = await parser.parse(b"", "empty.txt")

        assert result.status == "parsed"
        assert result.content == ""

    @pytest.mark.asyncio
    async def test_corrupted_pdf(self, parser):
        """Test parsing corrupted PDF."""
        with patch('pdfplumber.open', side_effect=Exception("Corrupted PDF")):
            result = await parser.parse(b"corrupted", "bad.pdf")

        assert result.status == "failed"
        assert result.error is not None

    @pytest.mark.asyncio
    async def test_corrupted_docx(self, parser):
        """Test parsing corrupted DOCX."""
        with patch('api.services.documents.DocxDocument', side_effect=Exception("Corrupted DOCX")):
            result = await parser.parse(b"corrupted", "bad.docx")

        assert result.status == "failed"

    @pytest.mark.asyncio
    async def test_corrupted_xlsx(self, parser):
        """Test parsing corrupted XLSX."""
        with patch('openpyxl.load_workbook', side_effect=Exception("Corrupted XLSX")):
            result = await parser.parse(b"corrupted", "bad.xlsx")

        assert result.status == "failed"

    @pytest.mark.asyncio
    async def test_file_extension_case_insensitive(self, parser, sample_text_bytes):
        """Test that file extension matching is case-insensitive."""
        result1 = await parser.parse(sample_text_bytes, "test.TXT")
        result2 = await parser.parse(sample_text_bytes, "test.Txt")

        assert result1.status == "parsed"
        assert result2.status == "parsed"

    @pytest.mark.asyncio
    async def test_filename_with_multiple_dots(self, parser, sample_text_bytes):
        """Test filename with multiple dots."""
        result = await parser.parse(sample_text_bytes, "my.document.v2.txt")

        assert result.status == "parsed"
        assert result.metadata["file_type"] == "txt"

    @pytest.mark.asyncio
    async def test_file_at_size_limit(self, parser):
        """Test file at exactly MAX_FILE_SIZE."""
        file_bytes = b"x" * MAX_FILE_SIZE

        result = await parser.parse(file_bytes, "large.txt")

        # Should succeed at exactly the limit
        assert result.status == "parsed"

    @pytest.mark.asyncio
    async def test_concurrent_parsing(self, parser, sample_text_bytes):
        """Test that multiple files can be parsed concurrently."""
        tasks = [
            parser.parse(sample_text_bytes, f"file{i}.txt")
            for i in range(5)
        ]

        results = await asyncio.gather(*tasks)

        assert len(results) == 5
        assert all(r.status == "parsed" for r in results)


# ============================================================================
# TEST MIME TYPE HANDLING
# ============================================================================

class TestMimeTypeHandling:
    """Tests for MIME type handling."""

    @pytest.mark.asyncio
    async def test_parse_with_mime_type_hint(self, parser, sample_text_bytes):
        """Test that mime_type parameter is accepted (even if not used)."""
        result = await parser.parse(
            sample_text_bytes,
            "document.txt",
            mime_type="text/plain"
        )

        assert result.status == "parsed"

    @pytest.mark.asyncio
    async def test_parse_relies_on_extension_not_mime(self, parser):
        """Test that parsing relies on file extension, not MIME type."""
        # Pass text content but claim it's a PDF via MIME
        result = await parser.parse(
            b"plain text",
            "document.pdf",  # Extension determines parser
            mime_type="text/plain"  # MIME type is ignored
        )

        # Should try PDF parsing (and likely fail), not text parsing
        # This tests that extension takes precedence
        assert result.metadata["file_type"] == "pdf"


# ============================================================================
# TEST CLEANUP AND RESOURCE MANAGEMENT
# ============================================================================

class TestResourceManagement:
    """Tests for temporary file cleanup and resource management."""

    def test_parser_does_not_have_persistent_temp_dir(self):
        """Test that parser does not create a persistent temporary directory."""
        with patch('api.services.documents.anthropic.Anthropic'):
            parser = DocumentParser()

            # Parser should not have temp_dir attribute (temp dirs are created per-operation)
            assert not hasattr(parser, 'temp_dir')

    @pytest.mark.asyncio
    async def test_temp_dir_cleaned_up_after_docx_conversion(self, parser):
        """Test that temp directory is cleaned up after DOC/ODT conversion."""
        mock_process = AsyncMock()
        mock_process.communicate = AsyncMock(return_value=(b"", b""))

        created_temp_dirs = []

        original_temp_dir = tempfile.TemporaryDirectory

        class TrackingTempDir:
            def __init__(self, *args, **kwargs):
                self._temp_dir = original_temp_dir(*args, **kwargs)
                created_temp_dirs.append(self._temp_dir.name)

            def __enter__(self):
                return self._temp_dir.__enter__()

            def __exit__(self, *args):
                return self._temp_dir.__exit__(*args)

        mock_doc = Mock()
        mock_doc.paragraphs = [Mock(text="Converted content")]
        mock_doc.tables = []

        with patch('tempfile.TemporaryDirectory', TrackingTempDir), \
             patch('asyncio.create_subprocess_exec', return_value=mock_process), \
             patch('builtins.open', mock_open(read_data=b"docx bytes")), \
             patch('os.path.exists', return_value=True), \
             patch('api.services.documents.DocxDocument', return_value=mock_doc):

            await parser.parse(b"doc content", "test.doc")

        # Verify temp dir was created and then cleaned up
        assert len(created_temp_dirs) == 1
        assert not os.path.exists(created_temp_dirs[0])

    @pytest.mark.asyncio
    async def test_temp_dir_cleaned_up_after_xlsx_conversion(self, parser):
        """Test that temp directory is cleaned up after XLS/ODS conversion."""
        mock_process = AsyncMock()
        mock_process.communicate = AsyncMock(return_value=(b"", b""))

        created_temp_dirs = []

        original_temp_dir = tempfile.TemporaryDirectory

        class TrackingTempDir:
            def __init__(self, *args, **kwargs):
                self._temp_dir = original_temp_dir(*args, **kwargs)
                created_temp_dirs.append(self._temp_dir.name)

            def __enter__(self):
                return self._temp_dir.__enter__()

            def __exit__(self, *args):
                return self._temp_dir.__exit__(*args)

        mock_cell = Mock()
        mock_cell.value = "Data"

        mock_sheet = Mock()
        mock_sheet.iter_rows.return_value = [[mock_cell]]

        mock_wb = Mock()
        mock_wb.sheetnames = ["Sheet1"]
        mock_wb.__getitem__ = Mock(return_value=mock_sheet)

        with patch('tempfile.TemporaryDirectory', TrackingTempDir), \
             patch('asyncio.create_subprocess_exec', return_value=mock_process), \
             patch('builtins.open', mock_open(read_data=b"xlsx bytes")), \
             patch('os.path.exists', return_value=True), \
             patch('openpyxl.load_workbook', return_value=mock_wb):

            await parser.parse(b"xls content", "test.xls")

        # Verify temp dir was created and then cleaned up
        assert len(created_temp_dirs) == 1
        assert not os.path.exists(created_temp_dirs[0])

    @pytest.mark.asyncio
    async def test_temp_dir_cleaned_up_after_pptx_conversion(self, parser):
        """Test that temp directory is cleaned up after PPT/ODP conversion."""
        mock_process = AsyncMock()
        mock_process.communicate = AsyncMock(return_value=(b"", b""))

        created_temp_dirs = []

        original_temp_dir = tempfile.TemporaryDirectory

        class TrackingTempDir:
            def __init__(self, *args, **kwargs):
                self._temp_dir = original_temp_dir(*args, **kwargs)
                created_temp_dirs.append(self._temp_dir.name)

            def __enter__(self):
                return self._temp_dir.__enter__()

            def __exit__(self, *args):
                return self._temp_dir.__exit__(*args)

        mock_shape = Mock()
        mock_shape.text = "Slide content"

        mock_slide = Mock()
        mock_slide.shapes = [mock_shape]

        mock_prs = Mock()
        mock_prs.slides = [mock_slide]

        with patch('tempfile.TemporaryDirectory', TrackingTempDir), \
             patch('asyncio.create_subprocess_exec', return_value=mock_process), \
             patch('builtins.open', mock_open(read_data=b"pptx bytes")), \
             patch('os.path.exists', return_value=True), \
             patch('api.services.documents.Presentation', return_value=mock_prs):

            await parser.parse(b"ppt content", "test.ppt")

        # Verify temp dir was created and then cleaned up
        assert len(created_temp_dirs) == 1
        assert not os.path.exists(created_temp_dirs[0])

    @pytest.mark.asyncio
    async def test_temp_dir_cleaned_up_after_rar_extraction(self, parser):
        """Test that temp directory is cleaned up after RAR extraction."""
        mock_process = AsyncMock()
        mock_process.communicate = AsyncMock(return_value=(b"", b""))

        created_temp_dirs = []

        original_temp_dir = tempfile.TemporaryDirectory

        class TrackingTempDir:
            def __init__(self, *args, **kwargs):
                self._temp_dir = original_temp_dir(*args, **kwargs)
                created_temp_dirs.append(self._temp_dir.name)

            def __enter__(self):
                return self._temp_dir.__enter__()

            def __exit__(self, *args):
                return self._temp_dir.__exit__(*args)

        with patch('tempfile.TemporaryDirectory', TrackingTempDir), \
             patch('asyncio.create_subprocess_exec', return_value=mock_process), \
             patch('builtins.open', mock_open()), \
             patch('os.makedirs'), \
             patch('os.walk', return_value=[]):

            await parser.parse(b"rar content", "archive.rar")

        # Verify temp dir was created and then cleaned up
        assert len(created_temp_dirs) == 1
        assert not os.path.exists(created_temp_dirs[0])

    @pytest.mark.asyncio
    async def test_temp_dir_cleaned_up_after_msg_parsing(self, parser):
        """Test that temp directory is cleaned up after MSG parsing."""
        extract_msg = pytest.importorskip('extract_msg')

        created_temp_dirs = []

        original_temp_dir = tempfile.TemporaryDirectory

        class TrackingTempDir:
            def __init__(self, *args, **kwargs):
                self._temp_dir = original_temp_dir(*args, **kwargs)
                created_temp_dirs.append(self._temp_dir.name)

            def __enter__(self):
                return self._temp_dir.__enter__()

            def __exit__(self, *args):
                return self._temp_dir.__exit__(*args)

        mock_msg = Mock()
        mock_msg.sender = "sender@test.com"
        mock_msg.to = "recipient@test.com"
        mock_msg.subject = "Test"
        mock_msg.date = "2024-01-01"
        mock_msg.body = "Body"

        with patch('tempfile.TemporaryDirectory', TrackingTempDir), \
             patch('builtins.open', mock_open()), \
             patch('extract_msg.Message', return_value=mock_msg):

            await parser.parse(b"msg content", "message.msg")

        # Verify temp dir was created and then cleaned up
        assert len(created_temp_dirs) == 1
        assert not os.path.exists(created_temp_dirs[0])

    @pytest.mark.asyncio
    async def test_temp_dir_cleaned_up_on_conversion_error(self, parser):
        """Test that temp directory is cleaned up even when conversion fails."""
        mock_process = AsyncMock()
        mock_process.communicate = AsyncMock(return_value=(b"", b""))

        created_temp_dirs = []

        original_temp_dir = tempfile.TemporaryDirectory

        class TrackingTempDir:
            def __init__(self, *args, **kwargs):
                self._temp_dir = original_temp_dir(*args, **kwargs)
                created_temp_dirs.append(self._temp_dir.name)

            def __enter__(self):
                return self._temp_dir.__enter__()

            def __exit__(self, *args):
                return self._temp_dir.__exit__(*args)

        with patch('tempfile.TemporaryDirectory', TrackingTempDir), \
             patch('asyncio.create_subprocess_exec', return_value=mock_process), \
             patch('builtins.open', mock_open()), \
             patch('os.path.exists', return_value=False):  # Simulate conversion failure

            result = await parser.parse(b"doc content", "test.doc")

        # Conversion should fail but temp dir should still be cleaned up
        assert result.status == "failed"
        assert len(created_temp_dirs) == 1
        assert not os.path.exists(created_temp_dirs[0])

    @pytest.mark.asyncio
    async def test_no_temp_files_leak_after_multiple_parses(self, parser):
        """Test that no temp files leak after parsing multiple documents."""
        # Get initial temp dir contents
        temp_base = tempfile.gettempdir()
        initial_dirs = set(os.listdir(temp_base))

        mock_process = AsyncMock()
        mock_process.communicate = AsyncMock(return_value=(b"", b""))

        mock_doc = Mock()
        mock_doc.paragraphs = [Mock(text="Content")]
        mock_doc.tables = []

        # Parse multiple documents
        for i in range(5):
            with patch('asyncio.create_subprocess_exec', return_value=mock_process), \
                 patch('builtins.open', mock_open(read_data=b"docx bytes")), \
                 patch('os.path.exists', return_value=True), \
                 patch('api.services.documents.DocxDocument', return_value=mock_doc):

                await parser.parse(b"doc content", f"test{i}.doc")

        # Check that no new temp directories remain
        final_dirs = set(os.listdir(temp_base))
        new_dirs = final_dirs - initial_dirs

        # Filter to only temp dirs that might be from our parser
        # (they should have been cleaned up)
        parser_temp_dirs = [d for d in new_dirs if d.startswith('tmp')]

        # All temp dirs should be cleaned up
        # (allowing some tolerance for other system processes)
        assert len(parser_temp_dirs) == 0, f"Leaked temp dirs: {parser_temp_dirs}"


# ============================================================================
# TEST INTEGRATION SCENARIOS
# ============================================================================

class TestIntegrationScenarios:
    """Integration tests for realistic scenarios."""

    @pytest.mark.asyncio
    async def test_parse_resume_pdf(self, parser):
        """Test parsing a typical resume PDF."""
        mock_page = Mock()
        mock_page.extract_text.return_value = """
        John Doe
        Software Engineer

        Experience:
        - Company A: Senior Developer (2020-2024)
        - Company B: Junior Developer (2018-2020)

        Skills: Python, JavaScript, SQL
        """
        mock_page.extract_tables.return_value = []

        mock_pdf = Mock()
        mock_pdf.pages = [mock_page]
        mock_pdf.__enter__ = Mock(return_value=mock_pdf)
        mock_pdf.__exit__ = Mock(return_value=False)

        with patch('pdfplumber.open', return_value=mock_pdf):
            result = await parser.parse(b"pdf bytes", "resume.pdf")

        assert result.status == "parsed"
        assert "John Doe" in result.content
        assert "Software Engineer" in result.content
        assert "Python" in result.content

    @pytest.mark.asyncio
    async def test_parse_job_application_zip(self, parser):
        """Test parsing a job application ZIP with multiple files."""
        zip_buffer = io.BytesIO()
        with zipfile.ZipFile(zip_buffer, 'w') as zf:
            zf.writestr("resume.txt", "John Doe - Software Engineer")
            zf.writestr("cover_letter.txt", "Dear Hiring Manager...")
            zf.writestr("references.txt", "Reference 1: Jane Smith")

        zip_bytes = zip_buffer.getvalue()

        result = await parser.parse(zip_bytes, "application.zip")

        assert result.status == "parsed"
        assert "resume.txt" in result.content
        assert "John Doe" in result.content
        assert "cover_letter.txt" in result.content

    @pytest.mark.asyncio
    async def test_parse_employee_data_xlsx(self, parser):
        """Test parsing employee data spreadsheet."""
        mock_cells_row1 = [
            Mock(value="Name"),
            Mock(value="Department"),
            Mock(value="Position")
        ]
        mock_cells_row2 = [
            Mock(value="John Doe"),
            Mock(value="IT"),
            Mock(value="Engineer")
        ]

        mock_sheet = Mock()
        mock_sheet.iter_rows.return_value = [mock_cells_row1, mock_cells_row2]

        mock_wb = Mock()
        mock_wb.sheetnames = ["Employees"]
        mock_wb.__getitem__ = Mock(return_value=mock_sheet)

        with patch('openpyxl.load_workbook', return_value=mock_wb):
            result = await parser.parse(b"xlsx content", "employees.xlsx")

        assert result.status == "parsed"
        assert "Employees" in result.content
        assert "John Doe" in result.content
        assert "IT" in result.content


# ============================================================================
# TEST ENCODING DETECTION AND GARBLED TEXT
# ============================================================================

class TestEncodingDetection:
    """Tests for encoding detection and garbled text handling."""

    @pytest.mark.asyncio
    async def test_parse_cp1251_russian_text(self, parser):
        """Test parsing Russian text in CP1251 encoding."""
        russian_text = "Привет, мир! Это тестовый документ на русском языке."
        cp1251_bytes = russian_text.encode('cp1251')

        result = await parser.parse(cp1251_bytes, "russian.txt")

        assert result.status == "parsed"
        assert "Привет" in result.content
        assert "русском" in result.content
        # Should have detected and used cp1251
        assert "cp1251" in result.metadata.get("encoding", "").lower() or "1251" in result.metadata.get("encoding", "")

    @pytest.mark.asyncio
    async def test_parse_koi8r_russian_text(self, parser):
        """Test parsing Russian text in KOI8-R encoding."""
        russian_text = "Добро пожаловать! Тестирование кодировки KOI8-R."
        koi8r_bytes = russian_text.encode('koi8-r')

        result = await parser.parse(koi8r_bytes, "koi8_russian.txt")

        assert result.status == "parsed"
        assert "Добро" in result.content or "пожаловать" in result.content

    @pytest.mark.asyncio
    async def test_detect_cjk_garbled_text(self, parser):
        """Test detection of CJK characters as garbled text indicator."""
        # Simulate garbled text: Russian CP1251 interpreted as UTF-8 produces CJK
        # This is a common encoding issue
        russian_text = "Тестовый документ"
        cp1251_bytes = russian_text.encode('cp1251')

        # The parser should detect CJK mixed with Cyrillic as garbled
        # and try alternative encodings
        result = await parser.parse(cp1251_bytes, "garbled.txt")

        assert result.status == "parsed"
        # Should NOT contain CJK characters
        content = result.content
        cjk_count = sum(1 for c in content if 0x4E00 <= ord(c) <= 0x9FFF)
        assert cjk_count == 0, f"Found {cjk_count} CJK characters in parsed content"

    @pytest.mark.asyncio
    async def test_has_garbled_text_detects_cjk_mixed_with_cyrillic(self, parser):
        """Test _has_garbled_text detects CJK mixed with Cyrillic."""
        # Text with mixed CJK and Cyrillic - definitely garbled
        mixed_text = "Привет 你好 world"
        assert parser._has_garbled_text(mixed_text) is True

    @pytest.mark.asyncio
    async def test_has_garbled_text_allows_pure_cyrillic(self, parser):
        """Test _has_garbled_text allows pure Cyrillic text."""
        cyrillic_text = "Чисто русский текст без проблем"
        assert parser._has_garbled_text(cyrillic_text) is False

    @pytest.mark.asyncio
    async def test_has_garbled_text_allows_pure_latin(self, parser):
        """Test _has_garbled_text allows pure Latin text."""
        latin_text = "Pure English text without issues"
        assert parser._has_garbled_text(latin_text) is False

    @pytest.mark.asyncio
    async def test_has_garbled_text_detects_high_cjk_ratio(self, parser):
        """Test _has_garbled_text detects text with high CJK ratio (likely wrong encoding)."""
        # Text with many CJK chars but no Cyrillic - suspicious for Russian HR bot
        cjk_heavy_text = "你好世界测试文档内容"
        # This should be flagged as suspicious in Russian context
        result = parser._has_garbled_text(cjk_heavy_text)
        # Note: depends on implementation - may or may not flag pure CJK
        assert isinstance(result, bool)

    @pytest.mark.asyncio
    async def test_quality_score_prefers_readable_text(self, parser):
        """Test that quality score prefers readable Cyrillic/Latin text."""
        readable_text = "Хороший читаемый текст с нормальными символами"
        garbled_text = "你好世界测试" + "\ufffd" * 10

        readable_score = parser._calculate_text_quality_score(readable_text)
        garbled_score = parser._calculate_text_quality_score(garbled_text)

        assert readable_score > garbled_score

    @pytest.mark.asyncio
    async def test_parse_html_with_cp1251_encoding(self, parser):
        """Test parsing HTML file with CP1251 encoding."""
        html_content = """<!DOCTYPE html>
<html>
<head>
    <meta charset="windows-1251">
    <title>Резюме</title>
</head>
<body>
    <h1>Иван Иванов</h1>
    <p>Программист Python</p>
</body>
</html>"""
        cp1251_bytes = html_content.encode('cp1251')

        result = await parser.parse(cp1251_bytes, "resume.html")

        assert result.status == "parsed"
        assert "Иван" in result.content or "Резюме" in result.content
        # Should NOT have garbled characters
        cjk_count = sum(1 for c in result.content if 0x4E00 <= ord(c) <= 0x9FFF)
        assert cjk_count == 0, f"HTML parsing produced {cjk_count} CJK characters"

    @pytest.mark.asyncio
    async def test_encoding_fallback_with_partial_status(self, parser):
        """Test that failed encoding detection results in partial status."""
        # Binary garbage that can't be properly decoded
        binary_garbage = bytes(range(256))

        result = await parser.parse(binary_garbage, "binary.txt")

        # Should either fail or mark as partial due to encoding issues
        assert result.status in ["parsed", "partial", "failed"]

    @pytest.mark.asyncio
    async def test_utf8_with_bom(self, parser):
        """Test parsing UTF-8 file with BOM."""
        bom = b'\xef\xbb\xbf'
        text = "Текст с BOM маркером"
        utf8_with_bom = bom + text.encode('utf-8')

        result = await parser.parse(utf8_with_bom, "bom.txt")

        assert result.status == "parsed"
        assert "Текст" in result.content

    @pytest.mark.asyncio
    async def test_mixed_language_document(self, parser):
        """Test parsing document with mixed Russian and English."""
        mixed_text = """
        Resume / Резюме

        Name: John Doe / Имя: Иван Иванов
        Position: Software Engineer / Должность: Программист

        Skills / Навыки:
        - Python, JavaScript
        - Разработка веб-приложений
        """

        result = await parser.parse(mixed_text.encode('utf-8'), "mixed.txt")

        assert result.status == "parsed"
        assert "Resume" in result.content
        assert "Резюме" in result.content
