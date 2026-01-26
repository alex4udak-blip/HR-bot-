import os
import io
import json
import tempfile
import subprocess
import asyncio
import zipfile
import chardet
from pathlib import Path
from typing import Dict, Any, Optional, Tuple
from datetime import datetime
import logging

import aiofiles
import pdfplumber
from docx import Document as DocxDocument
from pptx import Presentation
import openpyxl
import pandas as pd
from bs4 import BeautifulSoup
from striprtf.striprtf import rtf_to_text
from PIL import Image
import anthropic

from ..config import settings

logger = logging.getLogger(__name__)

# File size limits
MAX_FILE_SIZE = 20 * 1024 * 1024  # 20MB
MAX_PDF_PAGES = 50
MAX_ARCHIVE_FILES = 10
PARSE_TIMEOUT = 60

# Async Anthropic client for non-blocking OCR calls
_async_anthropic_client: anthropic.AsyncAnthropic | None = None


def get_async_anthropic_client() -> anthropic.AsyncAnthropic:
    """Get or create async Anthropic client (singleton pattern)."""
    global _async_anthropic_client
    if _async_anthropic_client is None:
        _async_anthropic_client = anthropic.AsyncAnthropic(api_key=settings.anthropic_api_key)
    return _async_anthropic_client


class DocumentParseResult:
    def __init__(
        self,
        content: str = "",
        status: str = "parsed",
        metadata: Optional[Dict[str, Any]] = None,
        error: Optional[str] = None
    ):
        self.content = content
        self.status = status  # "parsed", "partial", "failed"
        self.metadata = metadata or {}
        self.error = error


class DocumentParser:
    def __init__(self):
        # Use async client to avoid blocking the event loop
        # Client is lazily initialized via get_async_anthropic_client()
        pass

    async def parse(self, file_bytes: bytes, filename: str, mime_type: Optional[str] = None) -> DocumentParseResult:
        """Main entry point for document parsing."""
        if len(file_bytes) > MAX_FILE_SIZE:
            return DocumentParseResult(
                status="failed",
                error=f"File too large: {len(file_bytes) / 1024 / 1024:.1f}MB (max {MAX_FILE_SIZE / 1024 / 1024}MB)",
                metadata={"filename": filename, "file_size": len(file_bytes)}
            )

        ext = Path(filename).suffix.lower().lstrip('.')

        parsers = {
            # Documents
            'pdf': self._parse_pdf,
            'docx': self._parse_docx,
            'doc': self._convert_then_parse_docx,
            'odt': self._convert_then_parse_docx,
            'rtf': self._parse_rtf,
            'txt': self._parse_text,
            'md': self._parse_text,
            'html': self._parse_html,
            'htm': self._parse_html,

            # Spreadsheets
            'xlsx': self._parse_xlsx,
            'xls': self._convert_then_parse_xlsx,
            'csv': self._parse_csv,
            'ods': self._convert_then_parse_xlsx,

            # Presentations
            'pptx': self._parse_pptx,
            'ppt': self._convert_then_parse_pptx,
            'odp': self._convert_then_parse_pptx,

            # Images (OCR)
            'jpg': self._ocr_with_vision,
            'jpeg': self._ocr_with_vision,
            'png': self._ocr_with_vision,
            'gif': self._ocr_with_vision,
            'webp': self._ocr_with_vision,
            'bmp': self._ocr_with_vision,
            'tiff': self._ocr_with_vision,
            'tif': self._ocr_with_vision,
            'heic': self._convert_then_ocr,

            # Archives
            'zip': self._extract_and_parse_archive,
            'rar': self._extract_and_parse_rar,
            '7z': self._extract_and_parse_7z,

            # Other
            'json': self._parse_json,
            'xml': self._parse_xml,
            'yaml': self._parse_text,
            'yml': self._parse_text,
            'eml': self._parse_email,
            'msg': self._parse_outlook_msg,
        }

        parser = parsers.get(ext)
        if not parser:
            return DocumentParseResult(
                status="failed",
                error=f"Unsupported file type: {ext}",
                metadata={"filename": filename, "file_type": ext}
            )

        try:
            result = await asyncio.wait_for(
                parser(file_bytes, filename),
                timeout=PARSE_TIMEOUT
            )
            result.metadata["filename"] = filename
            result.metadata["file_type"] = ext
            result.metadata["file_size"] = len(file_bytes)
            result.metadata["parsed_at"] = datetime.utcnow().isoformat()
            return result
        except asyncio.TimeoutError:
            return DocumentParseResult(
                status="failed",
                error="Parsing timeout exceeded",
                metadata={"filename": filename, "file_type": ext}
            )
        except Exception as e:
            logger.debug(f"Error parsing {filename}: {e}")
            return DocumentParseResult(
                status="failed",
                error=str(e),
                metadata={"filename": filename, "file_type": ext}
            )

    # ========== PDF ==========
    async def _parse_pdf(self, file_bytes: bytes, filename: str) -> DocumentParseResult:
        loop = asyncio.get_event_loop()
        return await loop.run_in_executor(None, self._parse_pdf_sync, file_bytes, filename)

    def _parse_pdf_sync(self, file_bytes: bytes, filename: str) -> DocumentParseResult:
        text_parts = []
        tables_count = 0

        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            pages_count = len(pdf.pages)
            pages_to_parse = min(pages_count, MAX_PDF_PAGES)

            for i, page in enumerate(pdf.pages[:pages_to_parse]):
                page_text = page.extract_text() or ""
                if page_text.strip():
                    text_parts.append(f"--- Page {i + 1} ---\n{page_text}")

                # Extract tables
                tables = page.extract_tables()
                for table in tables:
                    tables_count += 1
                    if table:
                        table_text = self._table_to_text(table)
                        text_parts.append(f"[Table {tables_count}]\n{table_text}")

        content = "\n\n".join(text_parts)
        status = "parsed" if content.strip() else "partial"

        return DocumentParseResult(
            content=content,
            status=status,
            metadata={
                "pages_count": pages_count,
                "pages_parsed": pages_to_parse,
                "tables_count": tables_count,
                "truncated": pages_count > MAX_PDF_PAGES
            }
        )

    def _table_to_text(self, table: list) -> str:
        """Convert table to readable text format."""
        lines = []
        for row in table:
            if row:
                cells = [str(cell or "").strip() for cell in row]
                lines.append(" | ".join(cells))
        return "\n".join(lines)

    # ========== DOCX ==========
    async def _parse_docx(self, file_bytes: bytes, filename: str) -> DocumentParseResult:
        loop = asyncio.get_event_loop()
        return await loop.run_in_executor(None, self._parse_docx_sync, file_bytes)

    def _parse_docx_sync(self, file_bytes: bytes) -> DocumentParseResult:
        doc = DocxDocument(io.BytesIO(file_bytes))
        paragraphs = []
        tables_count = 0

        for para in doc.paragraphs:
            if para.text.strip():
                paragraphs.append(para.text)

        # Extract tables
        for table in doc.tables:
            tables_count += 1
            table_text = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                table_text.append(" | ".join(cells))
            paragraphs.append(f"[Table {tables_count}]\n" + "\n".join(table_text))

        content = "\n\n".join(paragraphs)

        return DocumentParseResult(
            content=content,
            status="parsed" if content.strip() else "partial",
            metadata={
                "paragraphs_count": len(doc.paragraphs),
                "tables_count": tables_count
            }
        )

    # ========== DOC/ODT (convert to DOCX) ==========
    async def _convert_then_parse_docx(self, file_bytes: bytes, filename: str) -> DocumentParseResult:
        ext = Path(filename).suffix

        with tempfile.TemporaryDirectory() as temp_dir:
            temp_input = os.path.join(temp_dir, f"input{ext}")

            async with aiofiles.open(temp_input, 'wb') as f:
                await f.write(file_bytes)

            # Convert using LibreOffice
            process = await asyncio.create_subprocess_exec(
                'libreoffice', '--headless', '--convert-to', 'docx',
                '--outdir', temp_dir, temp_input,
                stdout=asyncio.subprocess.PIPE,
                stderr=asyncio.subprocess.PIPE
            )
            await process.communicate()

            # Find the output file
            output_name = Path(temp_input).stem + ".docx"
            output_path = os.path.join(temp_dir, output_name)

            if os.path.exists(output_path):
                async with aiofiles.open(output_path, 'rb') as f:
                    docx_bytes = await f.read()
                result = await self._parse_docx(docx_bytes, filename)
                result.metadata["converted_from"] = ext
                return result
            else:
                return DocumentParseResult(
                    status="failed",
                    error="LibreOffice conversion failed"
                )
        # temp_dir is automatically cleaned up when exiting the context

    # ========== RTF ==========
    async def _parse_rtf(self, file_bytes: bytes, filename: str) -> DocumentParseResult:
        loop = asyncio.get_event_loop()
        return await loop.run_in_executor(None, self._parse_rtf_sync, file_bytes)

    def _parse_rtf_sync(self, file_bytes: bytes) -> DocumentParseResult:
        # Try different encodings
        for encoding in ['utf-8', 'cp1251', 'latin-1']:
            try:
                rtf_content = file_bytes.decode(encoding)
                text = rtf_to_text(rtf_content)
                return DocumentParseResult(
                    content=text,
                    status="parsed" if text.strip() else "partial"
                )
            except (UnicodeDecodeError, Exception) as e:
                logger.warning(f"RTF decode failed with {encoding}: {e}")
                continue

        return DocumentParseResult(
            status="failed",
            error="Could not decode RTF file"
        )

    # ========== Plain Text ==========
    async def _parse_text(self, file_bytes: bytes, filename: str) -> DocumentParseResult:
        """Parse text file with robust encoding detection.

        Tries encodings in order of likelihood for Russian/CIS text:
        1. Use chardet detection if confidence is high (>0.7)
        2. Try UTF-8 first (most common modern encoding)
        3. Try Windows-1251 (common for Russian Windows files)
        4. Try KOI8-R (legacy Russian encoding)
        5. Fall back to chardet result with errors='replace'
        """
        detected = chardet.detect(file_bytes)
        detected_encoding = detected.get('encoding', 'utf-8') or 'utf-8'
        confidence = detected.get('confidence', 0)

        # List of encodings to try in order
        encodings_to_try = []

        # If chardet is confident, try its suggestion first
        if confidence > 0.7:
            encodings_to_try.append(detected_encoding)

        # Common encodings for Russian/CIS text
        # ALWAYS try cp1251 early - it's very common for Russian files
        # and chardet often misdetects it as UTF-8
        encodings_to_try.extend(['utf-8', 'cp1251', 'koi8-r', 'iso-8859-5', 'cp866', 'utf-16', 'utf-16-le', 'utf-16-be'])

        # If chardet was not confident, add its suggestion at the end
        if confidence <= 0.7 and detected_encoding not in encodings_to_try:
            encodings_to_try.append(detected_encoding)

        # Remove duplicates while preserving order
        seen = set()
        encodings_to_try = [e for e in encodings_to_try if not (e in seen or seen.add(e))]

        best_text = None
        best_encoding = None
        best_score = -1

        for encoding in encodings_to_try:
            try:
                text = file_bytes.decode(encoding)
                # Check the entire text for quality, not just sample
                sample = text[:2000]

                # Skip if has replacement characters
                if '\ufffd' in sample:
                    continue

                # Skip if has garbled text (CJK mixed with Cyrillic, etc.)
                if self._has_garbled_text(sample):
                    continue

                # Calculate quality score based on readable characters
                score = self._calculate_text_quality_score(text)

                if score > best_score:
                    best_score = score
                    best_text = text
                    best_encoding = encoding

                # If we have a high-quality result, use it
                if score > 0.8:
                    break

            except (UnicodeDecodeError, LookupError) as e:
                logger.debug(f"Decode with {encoding} failed: {e}")
                continue

        # Use best result found
        if best_text is not None:
            return DocumentParseResult(
                content=best_text,
                status="parsed",
                metadata={"encoding": best_encoding, "chardet_encoding": detected_encoding, "chardet_confidence": confidence, "quality_score": best_score}
            )

        # Final fallback - use detected encoding with replacement
        text = file_bytes.decode(detected_encoding, errors='replace')
        used_encoding = f"{detected_encoding} (with replacements)"
        logger.warning(f"All encoding attempts failed for {filename}, using {detected_encoding} with replacements")

        return DocumentParseResult(
            content=text,
            status="partial",  # Mark as partial since encoding wasn't clean
            metadata={"encoding": used_encoding, "chardet_encoding": detected_encoding, "chardet_confidence": confidence, "quality_score": 0}
        )

    def _calculate_text_quality_score(self, text: str) -> float:
        """Calculate quality score for decoded text.

        Higher score = more likely correct encoding.
        Based on ratio of readable characters (Cyrillic, Latin, digits, punctuation).
        """
        if not text:
            return 0.0

        readable = 0
        total = len(text)

        for char in text:
            code = ord(char)
            # Readable characters: Cyrillic, Latin, digits, common punctuation, whitespace
            if (0x0400 <= code <= 0x04FF) or \
               (0x0041 <= code <= 0x007A) or \
               (0x0030 <= code <= 0x0039) or \
               char in ' \n\t\r.,;:!?-—–\'\"()[]{}«»„""\'\'@#$%^&*+=/<>\\|`~_' or \
               code in (9, 10, 13, 32):  # Whitespace
                readable += 1

        return readable / total if total > 0 else 0.0

    def _has_garbled_text(self, text: str) -> bool:
        """Check if text appears to be garbled (wrong encoding).

        Returns True if text has suspicious patterns like:
        - High ratio of non-printable characters
        - CJK characters mixed with Cyrillic/Latin (encoding mismatch)
        - Many consecutive question marks or squares
        - Unusual character sequences
        """
        if not text:
            return False

        # Count character types
        suspicious = 0
        cjk_count = 0
        cyrillic_count = 0
        latin_count = 0
        total = len(text)

        for char in text:
            code = ord(char)

            # Count Cyrillic characters (Russian)
            if 0x0400 <= code <= 0x04FF:
                cyrillic_count += 1
            # Count Latin characters
            elif (0x0041 <= code <= 0x007A):  # A-Z, a-z
                latin_count += 1
            # Count CJK characters (Chinese/Japanese/Korean)
            # These appear when CP1251/KOI8-R is wrongly decoded as UTF-8
            elif (0x4E00 <= code <= 0x9FFF) or \
                 (0x3400 <= code <= 0x4DBF) or \
                 (0x20000 <= code <= 0x2A6DF) or \
                 (0xF900 <= code <= 0xFAFF) or \
                 (0x3000 <= code <= 0x303F) or \
                 (0x30A0 <= code <= 0x30FF) or \
                 (0x3040 <= code <= 0x309F):  # Hiragana/Katakana
                cjk_count += 1

            # Suspicious if:
            # - Control chars (except newline, tab, carriage return)
            # - Private use area
            # - Replacement characters
            if (code < 32 and code not in (9, 10, 13)) or \
               (0xE000 <= code <= 0xF8FF) or \
               (0xFFF0 <= code <= 0xFFFF) or \
               code == 0xFFFD:  # Replacement character
                suspicious += 1

        if total == 0:
            return False

        # If more than 5% suspicious characters, probably wrong encoding
        if suspicious / total > 0.05:
            return True

        # CJK characters mixed with Cyrillic = wrong encoding
        # This happens when CP1251/KOI8-R Russian text is decoded as UTF-8
        if cjk_count > 0 and cyrillic_count > 0:
            return True

        # If we have CJK characters but no Cyrillic, and text seems like it should be Russian
        # (common in Russian HR bot context), it's likely garbled
        if cjk_count > total * 0.1 and cyrillic_count == 0 and latin_count < total * 0.3:
            return True

        return False

    # ========== HTML ==========
    async def _parse_html(self, file_bytes: bytes, filename: str) -> DocumentParseResult:
        text_result = await self._parse_text(file_bytes, filename)

        soup = BeautifulSoup(text_result.content, 'html.parser')

        # Remove script and style elements
        for script in soup(["script", "style"]):
            script.decompose()

        text = soup.get_text(separator="\n")
        # Clean up whitespace
        lines = [line.strip() for line in text.splitlines() if line.strip()]
        text = "\n".join(lines)

        return DocumentParseResult(
            content=text,
            status="parsed" if text.strip() else "partial",
            metadata={"title": soup.title.string if soup.title else None}
        )

    # ========== XLSX ==========
    async def _parse_xlsx(self, file_bytes: bytes, filename: str) -> DocumentParseResult:
        loop = asyncio.get_event_loop()
        return await loop.run_in_executor(None, self._parse_xlsx_sync, file_bytes)

    def _parse_xlsx_sync(self, file_bytes: bytes) -> DocumentParseResult:
        wb = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
        sheets_data = []

        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            rows = []
            for row in sheet.iter_rows(max_row=100):  # Limit rows
                cells = [str(cell.value or "") for cell in row]
                if any(cells):
                    rows.append(" | ".join(cells))

            if rows:
                sheets_data.append(f"=== Sheet: {sheet_name} ===\n" + "\n".join(rows))

        content = "\n\n".join(sheets_data)

        return DocumentParseResult(
            content=content,
            status="parsed" if content.strip() else "partial",
            metadata={
                "sheets": wb.sheetnames,
                "sheets_count": len(wb.sheetnames)
            }
        )

    # ========== XLS/ODS (convert to XLSX) ==========
    async def _convert_then_parse_xlsx(self, file_bytes: bytes, filename: str) -> DocumentParseResult:
        ext = Path(filename).suffix

        with tempfile.TemporaryDirectory() as temp_dir:
            temp_input = os.path.join(temp_dir, f"input{ext}")

            async with aiofiles.open(temp_input, 'wb') as f:
                await f.write(file_bytes)

            process = await asyncio.create_subprocess_exec(
                'libreoffice', '--headless', '--convert-to', 'xlsx',
                '--outdir', temp_dir, temp_input,
                stdout=asyncio.subprocess.PIPE,
                stderr=asyncio.subprocess.PIPE
            )
            await process.communicate()

            output_name = Path(temp_input).stem + ".xlsx"
            output_path = os.path.join(temp_dir, output_name)

            if os.path.exists(output_path):
                async with aiofiles.open(output_path, 'rb') as f:
                    xlsx_bytes = await f.read()
                result = await self._parse_xlsx(xlsx_bytes, filename)
                result.metadata["converted_from"] = ext
                return result
            else:
                return DocumentParseResult(status="failed", error="Conversion failed")
        # temp_dir is automatically cleaned up when exiting the context

    # ========== CSV ==========
    async def _parse_csv(self, file_bytes: bytes, filename: str) -> DocumentParseResult:
        text_result = await self._parse_text(file_bytes, filename)

        try:
            df = pd.read_csv(io.StringIO(text_result.content), nrows=100)
            content = df.to_string()
            return DocumentParseResult(
                content=content,
                status="parsed",
                metadata={
                    "columns": list(df.columns),
                    "rows_count": len(df)
                }
            )
        except Exception as e:
            return DocumentParseResult(
                content=text_result.content,
                status="partial",
                error=f"CSV parsing failed: {e}"
            )

    # ========== PPTX ==========
    async def _parse_pptx(self, file_bytes: bytes, filename: str) -> DocumentParseResult:
        loop = asyncio.get_event_loop()
        return await loop.run_in_executor(None, self._parse_pptx_sync, file_bytes)

    def _parse_pptx_sync(self, file_bytes: bytes) -> DocumentParseResult:
        prs = Presentation(io.BytesIO(file_bytes))
        slides_text = []

        for i, slide in enumerate(prs.slides):
            slide_content = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content.append(shape.text)

            if slide_content:
                slides_text.append(f"--- Slide {i + 1} ---\n" + "\n".join(slide_content))

        content = "\n\n".join(slides_text)

        return DocumentParseResult(
            content=content,
            status="parsed" if content.strip() else "partial",
            metadata={"slides_count": len(prs.slides)}
        )

    # ========== PPT/ODP (convert to PPTX) ==========
    async def _convert_then_parse_pptx(self, file_bytes: bytes, filename: str) -> DocumentParseResult:
        ext = Path(filename).suffix

        with tempfile.TemporaryDirectory() as temp_dir:
            temp_input = os.path.join(temp_dir, f"input{ext}")

            async with aiofiles.open(temp_input, 'wb') as f:
                await f.write(file_bytes)

            process = await asyncio.create_subprocess_exec(
                'libreoffice', '--headless', '--convert-to', 'pptx',
                '--outdir', temp_dir, temp_input,
                stdout=asyncio.subprocess.PIPE,
                stderr=asyncio.subprocess.PIPE
            )
            await process.communicate()

            output_name = Path(temp_input).stem + ".pptx"
            output_path = os.path.join(temp_dir, output_name)

            if os.path.exists(output_path):
                async with aiofiles.open(output_path, 'rb') as f:
                    pptx_bytes = await f.read()
                result = await self._parse_pptx(pptx_bytes, filename)
                result.metadata["converted_from"] = ext
                return result
            else:
                return DocumentParseResult(status="failed", error="Conversion failed")
        # temp_dir is automatically cleaned up when exiting the context

    # ========== Images (OCR with Claude Vision) ==========
    async def _ocr_with_vision(self, file_bytes: bytes, filename: str) -> DocumentParseResult:
        """Perform OCR using Claude Vision API (async, non-blocking)."""
        import base64

        # Determine media type
        ext = Path(filename).suffix.lower().lstrip('.')
        media_types = {
            'jpg': 'image/jpeg',
            'jpeg': 'image/jpeg',
            'png': 'image/png',
            'gif': 'image/gif',
            'webp': 'image/webp',
            'bmp': 'image/bmp',
            'tiff': 'image/tiff',
            'tif': 'image/tiff',
        }
        media_type = media_types.get(ext, 'image/png')

        # Encode image
        image_data = base64.standard_b64encode(file_bytes).decode('utf-8')

        try:
            client = get_async_anthropic_client()
            response = await client.messages.create(
                model=settings.claude_model,
                max_tokens=4096,
                messages=[{
                    "role": "user",
                    "content": [
                        {
                            "type": "image",
                            "source": {
                                "type": "base64",
                                "media_type": media_type,
                                "data": image_data,
                            }
                        },
                        {
                            "type": "text",
                            "text": "Extract all text from this image. If it's a document (resume, CV, certificate, etc.), preserve the structure. If there's no text, describe what you see briefly."
                        }
                    ]
                }]
            )

            content = response.content[0].text

            return DocumentParseResult(
                content=content,
                status="parsed",
                metadata={"ocr_method": "claude_vision"}
            )
        except Exception as e:
            return DocumentParseResult(
                status="failed",
                error=f"OCR failed: {e}"
            )

    # ========== HEIC (convert then OCR) ==========
    async def _convert_then_ocr(self, file_bytes: bytes, filename: str) -> DocumentParseResult:
        """Convert HEIC image to JPEG, then perform OCR."""
        loop = asyncio.get_event_loop()

        # Convert HEIC to JPEG in thread pool (CPU-bound)
        try:
            jpeg_bytes = await loop.run_in_executor(None, self._convert_heic_to_jpeg, file_bytes)
        except Exception as e:
            return DocumentParseResult(
                status="failed",
                error=f"HEIC conversion failed: {e}"
            )

        # Now OCR the converted JPEG (async API call)
        return await self._ocr_with_vision(jpeg_bytes, filename.replace('.heic', '.jpg'))

    def _convert_heic_to_jpeg(self, file_bytes: bytes) -> bytes:
        """Synchronous HEIC to JPEG conversion (for thread pool execution)."""
        import pillow_heif

        pillow_heif.register_heif_opener()
        img = Image.open(io.BytesIO(file_bytes))

        # Convert to JPEG
        output = io.BytesIO()
        img.convert('RGB').save(output, format='JPEG', quality=95)
        return output.getvalue()

    # ========== Archives ==========
    async def _extract_and_parse_archive(self, file_bytes: bytes, filename: str) -> DocumentParseResult:
        """Extract and parse ZIP archive contents asynchronously."""
        extracted_files = []
        all_content = []

        try:
            # Extract file list synchronously (fast, in-memory)
            with zipfile.ZipFile(io.BytesIO(file_bytes)) as zf:
                file_list = zf.namelist()[:MAX_ARCHIVE_FILES]

                # Read all files from archive first (sync, fast)
                files_to_parse = []
                for name in file_list:
                    if zf.getinfo(name).is_dir():
                        continue
                    try:
                        content = zf.read(name)
                        files_to_parse.append((name, content))
                    except Exception as e:
                        extracted_files.append({
                            "name": name,
                            "status": "failed",
                            "error": str(e)
                        })

            # Parse extracted files asynchronously
            for name, content in files_to_parse:
                try:
                    # Recursively parse using await (not asyncio.run!)
                    result = await self.parse(content, name)

                    # Add to extracted_files regardless of content
                    file_info = {
                        "name": name,
                        "status": result.status,
                        "size": len(content)
                    }
                    if result.error:
                        file_info["error"] = result.error
                    extracted_files.append(file_info)

                    # Only add to content if there's actual content
                    if result.content:
                        all_content.append(f"=== {name} ===\n{result.content}")
                except Exception as e:
                    extracted_files.append({
                        "name": name,
                        "status": "failed",
                        "error": str(e)
                    })

            return DocumentParseResult(
                content="\n\n".join(all_content),
                status="parsed" if all_content else "partial",
                metadata={
                    "archive_type": "zip",
                    "extracted_files": extracted_files,
                    "total_files": len(file_list)
                }
            )
        except Exception as e:
            return DocumentParseResult(
                status="failed",
                error=f"ZIP extraction failed: {e}"
            )

    async def _extract_and_parse_rar(self, file_bytes: bytes, filename: str) -> DocumentParseResult:
        with tempfile.TemporaryDirectory() as temp_dir:
            temp_rar = os.path.join(temp_dir, "archive.rar")
            extract_dir = os.path.join(temp_dir, "extracted")

            async with aiofiles.open(temp_rar, 'wb') as f:
                await f.write(file_bytes)

            os.makedirs(extract_dir, exist_ok=True)

            process = await asyncio.create_subprocess_exec(
                'unrar', 'x', '-y', temp_rar, extract_dir,
                stdout=asyncio.subprocess.PIPE,
                stderr=asyncio.subprocess.PIPE
            )
            await process.communicate()

            all_content = []
            extracted_files = []

            for root, _, files in os.walk(extract_dir):
                for name in files[:MAX_ARCHIVE_FILES]:
                    filepath = os.path.join(root, name)
                    try:
                        async with aiofiles.open(filepath, 'rb') as f:
                            content = await f.read()
                        result = await self.parse(content, name)
                        if result.content:
                            all_content.append(f"=== {name} ===\n{result.content}")
                            extracted_files.append({"name": name, "status": result.status})
                    except Exception as e:
                        extracted_files.append({"name": name, "status": "failed"})

            return DocumentParseResult(
                content="\n\n".join(all_content),
                status="parsed" if all_content else "partial",
                metadata={
                    "archive_type": "rar",
                    "extracted_files": extracted_files
                }
            )
        # temp_dir is automatically cleaned up when exiting the context

    async def _extract_and_parse_7z(self, file_bytes: bytes, filename: str) -> DocumentParseResult:
        try:
            import py7zr

            all_content = []
            extracted_files = []

            with py7zr.SevenZipFile(io.BytesIO(file_bytes), 'r') as archive:
                for name, bio in archive.readall().items():
                    if len(extracted_files) >= MAX_ARCHIVE_FILES:
                        break
                    try:
                        content = bio.read()
                        result = await self.parse(content, name)
                        if result.content:
                            all_content.append(f"=== {name} ===\n{result.content}")
                            extracted_files.append({"name": name, "status": result.status})
                    except Exception as e:
                        logger.warning(f"7z file extraction failed for {name}: {e}")
                        extracted_files.append({"name": name, "status": "failed"})

            return DocumentParseResult(
                content="\n\n".join(all_content),
                status="parsed" if all_content else "partial",
                metadata={
                    "archive_type": "7z",
                    "extracted_files": extracted_files
                }
            )
        except Exception as e:
            return DocumentParseResult(status="failed", error=f"7z extraction failed: {e}")

    # ========== JSON ==========
    async def _parse_json(self, file_bytes: bytes, filename: str) -> DocumentParseResult:
        text_result = await self._parse_text(file_bytes, filename)

        try:
            data = json.loads(text_result.content)
            formatted = json.dumps(data, indent=2, ensure_ascii=False)
            return DocumentParseResult(
                content=formatted,
                status="parsed",
                metadata={"json_type": type(data).__name__}
            )
        except json.JSONDecodeError as e:
            logger.warning(f"JSON parsing failed: {e}")
            return DocumentParseResult(
                content=text_result.content,
                status="partial",
                error="Invalid JSON"
            )

    # ========== XML ==========
    async def _parse_xml(self, file_bytes: bytes, filename: str) -> DocumentParseResult:
        text_result = await self._parse_text(file_bytes, filename)

        try:
            soup = BeautifulSoup(text_result.content, 'xml')
            text = soup.get_text(separator="\n")
            return DocumentParseResult(
                content=text,
                status="parsed"
            )
        except Exception as e:
            logger.warning(f"XML parsing failed: {e}")
            return DocumentParseResult(
                content=text_result.content,
                status="partial"
            )

    # ========== Email ==========
    async def _parse_email(self, file_bytes: bytes, filename: str) -> DocumentParseResult:
        import email
        from email import policy

        try:
            msg = email.message_from_bytes(file_bytes, policy=policy.default)

            parts = []
            parts.append(f"From: {msg.get('From', 'Unknown')}")
            parts.append(f"To: {msg.get('To', 'Unknown')}")
            parts.append(f"Subject: {msg.get('Subject', 'No Subject')}")
            parts.append(f"Date: {msg.get('Date', 'Unknown')}")
            parts.append("---")

            # Get body
            if msg.is_multipart():
                for part in msg.walk():
                    if part.get_content_type() == "text/plain":
                        body = part.get_content()
                        parts.append(body)
            else:
                body = msg.get_content()
                parts.append(body)

            return DocumentParseResult(
                content="\n".join(parts),
                status="parsed",
                metadata={
                    "from": msg.get('From'),
                    "to": msg.get('To'),
                    "subject": msg.get('Subject')
                }
            )
        except Exception as e:
            return DocumentParseResult(status="failed", error=f"Email parsing failed: {e}")

    # ========== Outlook MSG ==========
    async def _parse_outlook_msg(self, file_bytes: bytes, filename: str) -> DocumentParseResult:
        with tempfile.TemporaryDirectory() as temp_dir:
            temp_msg = os.path.join(temp_dir, "message.msg")

            try:
                async with aiofiles.open(temp_msg, 'wb') as f:
                    await f.write(file_bytes)

                import extract_msg
                msg = extract_msg.Message(temp_msg)

                parts = []
                parts.append(f"From: {msg.sender}")
                parts.append(f"To: {msg.to}")
                parts.append(f"Subject: {msg.subject}")
                parts.append(f"Date: {msg.date}")
                parts.append("---")
                parts.append(msg.body or "")

                return DocumentParseResult(
                    content="\n".join(parts),
                    status="parsed",
                    metadata={
                        "from": msg.sender,
                        "to": msg.to,
                        "subject": msg.subject
                    }
                )
            except Exception as e:
                return DocumentParseResult(status="failed", error=f"MSG parsing failed: {e}")
        # temp_dir is automatically cleaned up when exiting the context


# Singleton instance
document_parser = DocumentParser()
